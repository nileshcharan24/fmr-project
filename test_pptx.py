"""Phase 4 test — generates a PPTX without calling Gemini."""
import sys
from pathlib import Path
sys.path.insert(0, str(Path(__file__).parent))

from backend.services.pptx_editor import run_pipeline

result = run_pipeline(
    username="testuser",
    company_name="Boat Lifestyle",
    tier=1,
    clusters=["Music", "Dance", "Proshows"],
    banner_count=3,
    logo_path="",
    manager_name="Arjun Sharma",
    manager_designation="Media Manager",
    manager_phone="9876543210",
    manager_email="arjun@festember.com",
    portfolio_name="Sonic Partner",
    fest_deliverables="Banners at 3 clusters\nMC mentions\n2 social media posts",
    company_deliverables="Social media posts\nGoodies for winners",
    brand_event_description="Silent Beat Disco using Boat headphones on college grounds.",
    outreach_city="Bangalore",
    include_csr=False,
)

print("Pipeline completed.")
print(f"Output folder : {result['output_folder']}")
print(f"PPTX path     : {result['pptx_path']}")
print(f"File exists   : {Path(result['pptx_path']).exists()}")
print(f"File size     : {Path(result['pptx_path']).stat().st_size // 1024} KB")

# Quick sanity check: open and count slides
from pptx import Presentation
prs = Presentation(result['pptx_path'])
print(f"Slide count   : {len(prs.slides)}")
print()
print("Slide preview:")
for i, slide in enumerate(prs.slides, 1):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                texts.append(t[:60])
    print(f"  Slide {i:2d}: {' | '.join(texts[:2])}")
