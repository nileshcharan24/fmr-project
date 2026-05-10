"""
Low-level PPTX operations.

Content replacement is done by directly editing the ZIP archive's XML bytes —
the only reliable method when placeholder text is split across multiple <a:r>
runs by PowerPoint's XML serializer.

Structural operations (slide deletion, duplication) use python-pptx internals.
"""
import copy
import re
import zipfile
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches


# ── Load / Save ───────────────────────────────────────────────────────────────

def load(path: str) -> Presentation:
    return Presentation(path)


def save(prs: Presentation, path: str) -> None:
    prs.save(path)


# ── Slide deletion ────────────────────────────────────────────────────────────

def _delete_slide_1based(prs: Presentation, one_based: int) -> None:
    zero = one_based - 1
    if zero < 0 or zero >= len(prs.slides):
        return
    sld_id_lst = prs.slides._sldIdLst
    sld_id_elem = sld_id_lst[zero]
    rid = sld_id_elem.get(qn('r:id'))
    sld_id_lst.remove(sld_id_elem)
    prs.part.drop_rel(rid)


def delete_slides_by_1based_indices(prs: Presentation, indices: list) -> None:
    """Delete slides by 1-based index. Always deletes highest first."""
    for idx in sorted(set(indices), reverse=True):
        _delete_slide_1based(prs, idx)


# ── Slide duplication ─────────────────────────────────────────────────────────

def duplicate_slide(prs: Presentation, source_1based: int) -> int:
    """
    Duplicate slide at source_1based (1-based), insert immediately after it.
    Returns 1-based index of the new slide.
    NOTE: Image relationships are fixed post-save by copy_slide_image_rels().
    """
    src = prs.slides[source_1based - 1]
    layout = src.slide_layout
    new_slide = prs.slides.add_slide(layout)

    # Replace new slide's spTree contents with a deep copy of source's
    src_tree = src.shapes._spTree
    new_tree = new_slide.shapes._spTree
    for child in list(new_tree):
        new_tree.remove(child)
    for child in src_tree:
        new_tree.append(copy.deepcopy(child))

    # Copy background
    src_bg = src._element.find(qn('p:bg'))
    if src_bg is not None:
        existing_bg = new_slide._element.find(qn('p:bg'))
        if existing_bg is not None:
            new_slide._element.remove(existing_bg)
        new_slide._element.insert(0, copy.deepcopy(src_bg))

    # Move from end of list to right after source
    sld_id_lst = prs.slides._sldIdLst
    new_elem = sld_id_lst[-1]
    sld_id_lst.remove(new_elem)
    sld_id_lst.insert(source_1based, new_elem)  # 0-based insert = right after source

    return source_1based + 1


def duplicate_slide_n_times(prs: Presentation, source_1based: int, extra_copies: int) -> list:
    """
    Duplicate cluster slide `extra_copies` times in sequence.
    Returns list of all 1-based cluster slide indices (original + copies).
    """
    indices = [source_1based]
    current = source_1based
    for _ in range(extra_copies):
        new_idx = duplicate_slide(prs, current)
        indices.append(new_idx)
        current = new_idx
    return indices


def duplicate_slide_zip(pptx_path: str, source_1based: int, count: int) -> None:
    """
    Duplicate a slide N times entirely at the ZIP level.

    python-pptx's slide duplication mangles image relationships — the new
    slides' rels point to wrong images even when the spTree is deep-copied
    correctly. This bypasses python-pptx entirely:

    Each new slide gets a verbatim copy of the source's slide XML AND a
    verbatim copy of the source's .rels file (with notesSlide stripped to
    avoid shared notes). presentation.xml.rels, presentation.xml's
    sldIdLst, and [Content_Types].xml are updated so PowerPoint can find
    the new slides.
    """
    if count <= 0:
        return

    pptx_bytes = Path(pptx_path).read_bytes()
    buf = BytesIO(pptx_bytes)

    with zipfile.ZipFile(buf, 'r') as zin:
        names = list(zin.namelist())
        ordered = _get_slide_order(zin)

        if source_1based > len(ordered) or source_1based < 1:
            return

        src_slide_path = ordered[source_1based - 1]
        src_filename = src_slide_path.split('/')[-1]
        src_xml = zin.read(src_slide_path)

        src_rels_path = f'ppt/slides/_rels/{src_filename}.rels'
        if src_rels_path in names:
            src_rels_str = zin.read(src_rels_path).decode('utf-8')
            # Strip notesSlide ref so each copy doesn't share the source's notes.
            src_rels_str = re.sub(
                r'<Relationship[^>]*Type="[^"]*notesSlide[^"]*"[^>]*/>',
                '',
                src_rels_str,
            )
            src_rels_bytes = src_rels_str.encode('utf-8')
        else:
            src_rels_bytes = b''

        # Find next available slideN.xml number
        existing_slides = [n for n in names if re.match(r'^ppt/slides/slide\d+\.xml$', n)]
        max_num = max(
            (int(re.search(r'slide(\d+)\.xml', n).group(1)) for n in existing_slides),
            default=0,
        )

        pres_rels_xml = zin.read('ppt/_rels/presentation.xml.rels').decode('utf-8')
        pres_xml      = zin.read('ppt/presentation.xml').decode('utf-8')
        ct_xml        = zin.read('[Content_Types].xml').decode('utf-8')

        # Locate source slide's rId in presentation.xml.rels
        src_pres_rid_match = re.search(
            r'<Relationship Id="(rId\d+)"[^>]*'
            r'Type="[^"]*relationships/slide"[^>]*'
            r'Target="(?:\.\./)?slides/' + re.escape(src_filename) + r'"',
            pres_rels_xml,
        )
        if not src_pres_rid_match:
            return
        src_pres_rid = src_pres_rid_match.group(1)

        max_rid = max(
            (int(r) for r in re.findall(r'Id="rId(\d+)"', pres_rels_xml)),
            default=0,
        )
        max_sld_id = max(
            (int(s) for s in re.findall(r'<p:sldId[^>]*id="(\d+)"', pres_xml)),
            default=255,
        )

        new_slides_data = []
        new_pres_rels = []
        new_sld_elems_str = ''
        new_overrides = []

        for i in range(count):
            new_num         = max_num + i + 1
            new_filename    = f'slide{new_num}.xml'
            new_slide_path  = f'ppt/slides/{new_filename}'
            new_rels_path   = f'ppt/slides/_rels/{new_filename}.rels'
            new_pres_rid    = f'rId{max_rid + i + 1}'
            new_sld_id      = max_sld_id + i + 1

            new_pres_rels.append(
                f'<Relationship Id="{new_pres_rid}" '
                f'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" '
                f'Target="slides/{new_filename}"/>'
            )
            # Append in order so they appear source, copy1, copy2, ... after insert
            new_sld_elems_str += f'<p:sldId id="{new_sld_id}" r:id="{new_pres_rid}"/>'
            new_overrides.append(
                f'<Override PartName="/ppt/slides/{new_filename}" '
                f'ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>'
            )
            new_slides_data.append((new_slide_path, new_rels_path, src_xml, src_rels_bytes))

        # Apply edits to the package metadata files
        pres_rels_xml = pres_rels_xml.replace(
            '</Relationships>',
            ''.join(new_pres_rels) + '</Relationships>',
        )
        # Insert new sldIds right after the source's sldId (single substitution = forward order)
        src_sld_pattern = re.compile(
            r'(<p:sldId[^/]*r:id="' + re.escape(src_pres_rid) + r'"\s*/>)'
        )
        pres_xml = src_sld_pattern.sub(r'\1' + new_sld_elems_str, pres_xml, count=1)

        for ov in new_overrides:
            if ov not in ct_xml:
                ct_xml = ct_xml.replace('</Types>', ov + '</Types>')

        # Write the modified zip
        out_buf = BytesIO()
        with zipfile.ZipFile(out_buf, 'w', zipfile.ZIP_DEFLATED) as zout:
            for name in names:
                data = zin.read(name)
                if name == 'ppt/_rels/presentation.xml.rels':
                    data = pres_rels_xml.encode('utf-8')
                elif name == 'ppt/presentation.xml':
                    data = pres_xml.encode('utf-8')
                elif name == '[Content_Types].xml':
                    data = ct_xml.encode('utf-8')
                zout.writestr(name, data)
            for new_slide_path, new_rels_path, slide_xml, rels_xml in new_slides_data:
                zout.writestr(new_slide_path, slide_xml)
                if rels_xml:
                    zout.writestr(new_rels_path, rels_xml)

    Path(pptx_path).write_bytes(out_buf.getvalue())


def cleanup_orphaned_slides(pptx_path: str) -> None:
    """
    Remove slide files and Content_Types overrides that are no longer
    referenced by presentation.xml.rels. python-pptx's slide deletion drops
    the relationship but leaves orphans behind — PowerPoint flags these as
    'unreadable content' on open and prompts to repair.

    Handles two kinds of orphan:
      1. Slide XML file present in the zip but not in rels → remove file + rels
         + Content_Types Override.
      2. Slide XML file already gone from the zip but Content_Types Override
         still references it → remove the Override.
    """
    pptx_bytes = Path(pptx_path).read_bytes()
    buf = BytesIO(pptx_bytes)

    with zipfile.ZipFile(buf, 'r') as zin:
        names = list(zin.namelist())
        referenced = set(_get_slide_order(zin))

        all_slide_files = {
            n for n in names
            if re.match(r'^ppt/slides/slide\d+\.xml$', n)
        }
        orphan_files = all_slide_files - referenced

        # Detect stale Content_Types overrides for slides not in rels
        ct_xml = zin.read('[Content_Types].xml').decode('utf-8')
        ct_slide_paths = re.findall(
            r'<Override\s+PartName="(/ppt/slides/slide\d+\.xml)"[^>]*/>',
            ct_xml,
        )
        stale_ct_paths = {
            p.lstrip('/') for p in ct_slide_paths if p.lstrip('/') not in referenced
        }

        if not orphan_files and not stale_ct_paths:
            return

        # Files to drop from the zip (orphan slide XML + matching .rels)
        skip = set(orphan_files)
        for o in orphan_files:
            slide_file = o.split('/')[-1]
            skip.add(f'ppt/slides/_rels/{slide_file}.rels')

        # All slide paths that need their Content_Types Override removed
        ct_paths_to_clear = orphan_files | stale_ct_paths

        out_buf = BytesIO()
        with zipfile.ZipFile(out_buf, 'w', zipfile.ZIP_DEFLATED) as zout:
            for name in names:
                if name in skip:
                    continue
                data = zin.read(name)
                if name == '[Content_Types].xml':
                    text = data.decode('utf-8')
                    for path in ct_paths_to_clear:
                        slide_file = path.split('/')[-1]
                        text = re.sub(
                            r'<Override\s+PartName="/ppt/slides/'
                            + re.escape(slide_file) + r'"[^>]*/>',
                            '',
                            text,
                        )
                    data = text.encode('utf-8')
                zout.writestr(name, data)

    Path(pptx_path).write_bytes(out_buf.getvalue())


def copy_slide_image_rels(pptx_path: str, source_1based: int, copy_count: int) -> None:
    """
    Fix missing image relationships in duplicated slides.

    When python-pptx duplicates a slide, the shape XML correctly references
    rId values from the source slide, but the new slide's .rels file only gets
    the layout relationship — not the image relationships. This function copies
    image rels from the source slide to each copy slide.

    Call this AFTER duplicate_slide_n_times() + save().
    """
    if copy_count <= 0:
        return

    pptx_bytes = Path(pptx_path).read_bytes()
    buf = BytesIO(pptx_bytes)

    with zipfile.ZipFile(buf, 'r') as zin:
        ordered_slides = _get_slide_order(zin)

        if source_1based > len(ordered_slides):
            return

        # Source slide XML path (e.g., 'ppt/slides/slide12.xml')
        src_slide_path = ordered_slides[source_1based - 1]
        # Rels file: 'ppt/slides/_rels/slideN.xml.rels'
        slide_file = src_slide_path.split('/')[-1]  # e.g., 'slide12.xml'
        src_rels_name = f'ppt/slides/_rels/{slide_file}.rels'

        if src_rels_name not in zin.namelist():
            return

        src_rels_xml = zin.read(src_rels_name).decode('utf-8')

        # Extract only image relationships from source rels
        image_rel_pattern = re.compile(
            r'<Relationship[^>]*Type="[^"]*relationships/image[^"]*"[^>]*/?>',
            re.IGNORECASE
        )
        image_rels = image_rel_pattern.findall(src_rels_xml)
        if not image_rels:
            return

        # Identify copy slides by their 1-based positions
        copy_rels_names = set()
        for i in range(copy_count):
            copy_idx = source_1based + i  # positions source+1, source+2, ...
            if copy_idx < len(ordered_slides):
                copy_slide_file = ordered_slides[copy_idx].split('/')[-1]
                copy_rels_names.add(f'ppt/slides/_rels/{copy_slide_file}.rels')

        out_buf = BytesIO()
        with zipfile.ZipFile(out_buf, 'w', zipfile.ZIP_DEFLATED) as zout:
            for name in zin.namelist():
                data = zin.read(name)
                if name in copy_rels_names:
                    rels_text = data.decode('utf-8')
                    # Inject image rels before closing tag (avoid duplicate rIds)
                    existing_ids = set(re.findall(r'Id="([^"]+)"', rels_text))
                    new_rels = []
                    for rel_str in image_rels:
                        rid_match = re.search(r'Id="([^"]+)"', rel_str)
                        if rid_match and rid_match.group(1) not in existing_ids:
                            new_rels.append(rel_str)
                            existing_ids.add(rid_match.group(1))
                    if new_rels:
                        rels_text = rels_text.replace(
                            '</Relationships>',
                            '\n'.join(new_rels) + '\n</Relationships>'
                        )
                    data = rels_text.encode('utf-8')
                zout.writestr(name, data)

    Path(pptx_path).write_bytes(out_buf.getvalue())


# ── Content replacement via ZIP byte editing ──────────────────────────────────

def _xml_escape(value: str) -> str:
    return (
        str(value)
        .replace('&', '&amp;')
        .replace('<', '&lt;')
        .replace('>', '&gt;')
        .replace('"', '&quot;')
    )


def _build_replacements_bytes(replacements: dict) -> list:
    """
    Pre-compile each replacement as (pattern, replacement_string).

    Placeholders in PPTX XML are HTML-entity encoded:
      <Company Name>  →  stored as  &lt;Company Name&gt;
    So we search for the &lt;...&gt; form, not literal angle brackets.

    The replacement value is plain text (not XML-escaped) because it goes
    inside an existing <a:t> text node — the surrounding XML is untouched.
    """
    compiled = []
    for key, value in replacements.items():
        escaped_key = re.escape(key)
        # Allow any whitespace (including newlines) between words in the key,
        # since PowerPoint sometimes stores <company\nname> with a literal newline.
        flexible_key = escaped_key.replace('\\ ', r'\s+').replace(' ', r'\s+')
        pattern = re.compile(
            r'&lt;' + flexible_key + r'[\s\w]*&gt;',
            re.IGNORECASE | re.DOTALL,
        )
        safe_value = _xml_escape(value)
        compiled.append((pattern, safe_value))
    return compiled


def _apply_split_run_replacements(text: str, replacements: dict) -> str:
    """
    Handle placeholders split across <a:t> nodes in three ways:

    Case A: &lt; alone in one <a:t>, then key&gt; in the next.
      e.g. <a:t>&lt;</a:t>...<a:t>company name&gt;</a:t>

    Case B: &lt;key_part1 in one <a:t>, then key_rest&gt; in the next (same paragraph).
      e.g. <a:t>&lt;company</a:t>...<a:t>name&gt;</a:t>

    Case C: same as Case B but split ACROSS paragraph boundaries (</a:p><a:p>).
      e.g. <a:t>&lt;company</a:t></a:p><a:p>...<a:t>name&gt;</a:t>
    """
    for key, value in replacements.items():
        safe_value = _xml_escape(value)
        _v = safe_value

        # Case A: &lt; alone in one node, key&gt; in next
        escaped_key = re.escape(key)
        pattern_a = re.compile(
            r'(&lt;)(</a:t>(?:(?!</a:p>).)*?<a:t>)'
            + escaped_key + r'[\s\w]*&gt;',
            re.IGNORECASE | re.DOTALL,
        )
        text = pattern_a.sub(lambda m, v=_v: v + m.group(2), text)

        # Case B & C: key split across multiple <a:t> nodes
        words = key.split(' ')
        if len(words) > 1:
            for split_at in range(1, len(words)):
                first = re.escape(' '.join(words[:split_at]))
                rest  = re.escape(' '.join(words[split_at:]))

                # Case B: same paragraph
                pattern_b = re.compile(
                    r'&lt;' + first + r'(</a:t>(?:(?!</a:p>).)*?<a:t>)'
                    + rest + r'[\s\w]*&gt;',
                    re.IGNORECASE | re.DOTALL,
                )
                text = pattern_b.sub(lambda m, v=_v: v + m.group(1), text)

                # Case C: across paragraph boundary
                pattern_c = re.compile(
                    r'&lt;' + first + r'(</a:t>.*?</a:p>\s*<a:p>.*?<a:t>)'
                    + rest + r'[\s\w]*&gt;',
                    re.IGNORECASE | re.DOTALL,
                )
                text = pattern_c.sub(lambda m, v=_v: v + m.group(1), text)

    return text


def _get_slide_order(zin: zipfile.ZipFile) -> list:
    """
    Read ppt/_rels/presentation.xml.rels to get slide XML filenames in
    presentation order. Returns list like ['ppt/slides/slide3.xml', ...].
    """
    rels_xml = zin.read('ppt/_rels/presentation.xml.rels').decode('utf-8')
    # Find all slide relationships in document order
    entries = re.findall(
        r'<Relationship[^>]+Type="[^"]*slide"[^>]+Target="([^"]+)"',
        rels_xml
    )
    # Targets are like '../slides/slide3.xml' or 'slides/slide3.xml'
    ordered = []
    for t in entries:
        # Normalise to 'ppt/slides/slideN.xml'
        name = t.replace('../', 'ppt/')
        if not name.startswith('ppt/'):
            name = 'ppt/' + name
        ordered.append(name)
    return ordered


def replace_in_slides_matching(
    pptx_path: str,
    match_text: str,
    replacements: dict,
    max_slides: int = None,
) -> None:
    """
    Apply replacements only to slides whose raw XML contains match_text.
    max_slides: if set, stop after replacing in this many slides (used to
                fill cluster copies one at a time).
    """
    compiled = _build_replacements_bytes(replacements)
    pptx_bytes = Path(pptx_path).read_bytes()
    buf = BytesIO(pptx_bytes)

    with zipfile.ZipFile(buf, 'r') as zin:
        ordered_slides = _get_slide_order(zin)  # presentation order
        # Find which slide files to replace (in presentation order)
        to_replace = set()
        for name in ordered_slides:
            if max_slides is not None and len(to_replace) >= max_slides:
                break
            text = zin.read(name).decode('utf-8')
            if match_text in text:
                to_replace.add(name)

        out_buf = BytesIO()
        with zipfile.ZipFile(out_buf, 'w', zipfile.ZIP_DEFLATED) as zout:
            for name in zin.namelist():
                data = zin.read(name)
                if name in to_replace:
                    text = data.decode('utf-8')
                    for pattern, safe_value in compiled:
                        text = pattern.sub(safe_value, text)
                    text = _apply_split_run_replacements(text, replacements)
                    data = text.encode('utf-8')
                zout.writestr(name, data)

    Path(pptx_path).write_bytes(out_buf.getvalue())


def replace_in_pptx_file(pptx_path: str, replacements: dict, slide_indices: list = None) -> None:
    """
    Open the PPTX zip, apply regex replacements to slide XML bytes in-place,
    and write back. Uses presentation order (not filename order) for indices.

    pptx_path    : path to the .pptx file (modified in place)
    replacements : { placeholder_text: replacement_value }
    slide_indices: 0-based indices into presentation slide order (None = all)
    """
    compiled = _build_replacements_bytes(replacements)
    pptx_bytes = Path(pptx_path).read_bytes()
    buf = BytesIO(pptx_bytes)

    with zipfile.ZipFile(buf, 'r') as zin:
        ordered_slides = _get_slide_order(zin)

        if slide_indices is not None:
            target_slides = {ordered_slides[i] for i in slide_indices if i < len(ordered_slides)}
        else:
            target_slides = set(ordered_slides)

        out_buf = BytesIO()
        with zipfile.ZipFile(out_buf, 'w', zipfile.ZIP_DEFLATED) as zout:
            for name in zin.namelist():
                data = zin.read(name)
                if name in target_slides:
                    text = data.decode('utf-8')
                    for pattern, safe_value in compiled:
                        text = pattern.sub(safe_value, text)
                    text = _apply_split_run_replacements(text, replacements)
                    data = text.encode('utf-8')
                zout.writestr(name, data)

    Path(pptx_path).write_bytes(out_buf.getvalue())


# ── Company deliverables structured slide fill ────────────────────────────────

# XML run properties for description text (light gray)
_DESC_RPR = (
    '<a:rPr lang="en-US" sz="2700">'
    '<a:solidFill><a:srgbClr val="ECECEC"/></a:solidFill>'
    '<a:latin typeface="Arial"/><a:ea typeface="Arial"/>'
    '<a:cs typeface="Arial"/><a:sym typeface="Arial"/>'
    '</a:rPr>'
)
_DESC_PARA_FMT = (
    '<a:p>'
    '<a:pPr marL="12700" lvl="0" indent="0" algn="l" rtl="0">'
    '<a:lnSpc><a:spcPct val="100000"/></a:lnSpc>'
    '<a:spcBef><a:spcPts val="0"/></a:spcBef>'
    '<a:spcAft><a:spcPts val="0"/></a:spcAft>'
    '<a:buNone/>'
    '</a:pPr>'
    '<a:r>{rpr}<a:t>{text}</a:t></a:r>'
    '<a:endParaRPr sz="2700">'
    '<a:latin typeface="Arial"/><a:ea typeface="Arial"/>'
    '<a:cs typeface="Arial"/><a:sym typeface="Arial"/>'
    '</a:endParaRPr>'
    '</a:p>'
)

# XML run properties for heading text (orange)
_HEAD_RPR = (
    '<a:rPr lang="en-US" sz="2700">'
    '<a:solidFill><a:srgbClr val="D89400"/></a:solidFill>'
    '<a:latin typeface="Arial"/><a:ea typeface="Arial"/>'
    '<a:cs typeface="Arial"/><a:sym typeface="Arial"/>'
    '</a:rPr>'
)


def _build_desc_xml(items_text: str) -> str:
    """
    Convert a newline-separated bullet list to PPTX paragraph XML.
    Each line becomes one <a:p> paragraph in the description style.
    """
    lines = [l.strip() for l in items_text.split('\n') if l.strip()]
    if not lines:
        return _DESC_PARA_FMT.format(rpr=_DESC_RPR, text='')
    paras = []
    for line in lines:
        paras.append(_DESC_PARA_FMT.format(rpr=_DESC_RPR, text=_xml_escape(line)))
    return '\n'.join(paras)


def fill_company_deliverables_slide(
    pptx_path: str,
    slide_position_1based: int,
    sections: list,
) -> None:
    """
    Fill the structured company deliverables slide.

    sections: list of (heading, items_text) tuples. Up to 4 entries.
      - heading: replaces 'SOCIAL MEDIA MARKETING', 'GOODIES', 'HEADING' in order
      - items_text: newline-separated bullet points replacing 'Description...'

    The slide structure (fixed in template):
      heading1 (orange, in title placeholder) — SOCIAL MEDIA MARKETING
      description1 (gray)                     — Description...
      heading2 (orange)                        — GOODIES
      description2 (gray)                     — Description...
      heading3 (orange)                        — HEADING
      description3 (gray)                     — Description...
      heading4 (orange)                        — HEADING
      description4 (gray)                     — Description...
    """
    pptx_bytes = Path(pptx_path).read_bytes()
    buf = BytesIO(pptx_bytes)

    with zipfile.ZipFile(buf, 'r') as zin:
        ordered_slides = _get_slide_order(zin)
        if slide_position_1based > len(ordered_slides):
            Path(pptx_path).write_bytes(pptx_bytes)
            return

        target_slide = ordered_slides[slide_position_1based - 1]

        out_buf = BytesIO()
        with zipfile.ZipFile(out_buf, 'w', zipfile.ZIP_DEFLATED) as zout:
            for name in zin.namelist():
                data = zin.read(name)
                if name == target_slide:
                    xml = data.decode('utf-8')
                    xml = _fill_deliverables_xml(xml, sections)
                    data = xml.encode('utf-8')
                zout.writestr(name, data)

    Path(pptx_path).write_bytes(out_buf.getvalue())


def _fill_deliverables_xml(xml: str, sections: list) -> str:
    """
    In-memory XML transformation for the deliverables slide.
    Replaces headings and Description... paragraphs sequentially.
    """
    fixed_headings = ['SOCIAL MEDIA MARKETING', 'GOODIES', 'HEADING', 'HEADING']

    for i, (heading, items_text) in enumerate(sections[:4]):
        old_head = fixed_headings[i]

        # Replace heading text (only first remaining occurrence of old_head)
        # Apply upper() BEFORE xml-escape so &amp; doesn't become &AMP;
        if heading and heading.upper() != old_head:
            xml = xml.replace(f'>{old_head}<', f'>{_xml_escape(heading.upper())}<', 1)

        # Replace the next 'Description...' paragraph with multi-line content
        desc_para_pattern = re.compile(
            r'<a:p>(?:(?!</a:p>).)*?<a:t>Description\.\.\.</a:t>.*?</a:p>',
            re.DOTALL
        )
        replacement_xml = _build_desc_xml(items_text)
        xml = desc_para_pattern.sub(replacement_xml, xml, count=1)

    return xml


def delete_slide_at_position(pptx_path: str, position_1based: int) -> None:
    """Delete a slide by its 1-based position in the current presentation order."""
    prs = load(pptx_path)
    delete_slides_by_1based_indices(prs, [position_1based])
    save(prs, pptx_path)


def find_slides_by_content(pptx_path: str, match_text: str) -> list:
    """
    Find all slides whose XML contains match_text. Returns a list of 1-based
    positions in presentation order. More robust than fixed indices when slides
    have been deleted or duplicated.
    """
    pptx_bytes = Path(pptx_path).read_bytes()
    buf = BytesIO(pptx_bytes)
    positions = []
    with zipfile.ZipFile(buf, 'r') as zin:
        ordered_slides = _get_slide_order(zin)
        for i, slide_name in enumerate(ordered_slides, start=1):
            try:
                text = zin.read(slide_name).decode('utf-8')
            except KeyError:
                continue
            if match_text in text:
                positions.append(i)
    return positions


# ── Logo insertion ────────────────────────────────────────────────────────────

def insert_logo(prs: Presentation, slide_1based: int, logo_path: str,
                left_in: float, top_in: float, width_in: float) -> None:
    slide = prs.slides[slide_1based - 1]
    slide.shapes.add_picture(
        logo_path,
        Inches(left_in),
        Inches(top_in),
        Inches(width_in),
    )
